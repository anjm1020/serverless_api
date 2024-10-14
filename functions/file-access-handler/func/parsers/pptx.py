from pptx import Presentation

from func.parsers.document_processor import DocumentProcessor


class PPTXProcessor(DocumentProcessor):
    def process(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
